from pptx import Presentation

prs = Presentation('template.pptx')

slide_layout = prs.slide_layouts[1]  # assuming you want the first one
slide = prs.slides.add_slide(slide_layout)

Presentation_Title = slide.placeholders[0]
Presentation_Subtitle = slide.placeholders[1]
Presentation_Title.text = 'This Is a Test'
Presentation_Subtitle.text = 'Is This Working?'

prs.save('SlideLayoutImportTest.pptx')
